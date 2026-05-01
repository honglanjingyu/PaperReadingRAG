# app/service/core/deepdoc/parser/ppt_parser.py

from pptx import Presentation


class RAGFlowPptParser:
    """PowerPoint 文档解析器"""

    def __call__(self, fnm, from_page, to_page, callback=None):
        ppt = Presentation(fnm) if isinstance(fnm, str) else Presentation(BytesIO(fnm))
        txts = []

        for i, slide in enumerate(ppt.slides):
            if i < from_page or i >= to_page:
                continue
            texts = []
            for shape in sorted(slide.shapes, key=lambda x: (x.top // 10, x.left)):
                txt = self.__extract(shape)
                if txt:
                    texts.append(txt)
            txts.append("\n".join(texts))
        return txts

    def __extract(self, shape):
        """提取形状中的文本"""
        if shape.shape_type == 19:  # 表格
            tb = shape.table
            rows = []
            for i in range(1, len(tb.rows)):
                rows.append("; ".join([tb.cell(0, j).text + ": " + tb.cell(i, j).text for j in range(len(tb.columns))]))
            return "\n".join(rows)

        if shape.has_text_frame:
            return shape.text_frame.text

        if shape.shape_type == 6:  # 组
            texts = []
            for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                t = self.__extract(p)
                if t:
                    texts.append(t)
            return "\n".join(texts)